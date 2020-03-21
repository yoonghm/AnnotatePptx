# MIT License
#
# Copyright (c) 2020 YoongHM
#
# Permission is hereby granted, free of charge, to any person obtaining a copy
# of this software and associated documentation files (the "Software"), to deal
# in the Software without restriction, including without limitation the rights
# to use, copy, modify, merge, publish, distribute, sublicense, and/or sell
# copies of the Software, and to permit persons to whom the Software is
# furnished to do so, subject to the following conditions:
#
# The above copyright notice and this permission notice shall be included in all
# copies or substantial portions of the Software.
#
# THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR
# IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY,
# FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE
# AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER
# LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM,
# OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE
# SOFTWARE.

# AnnotatePptx.py
# Version 1.0
#

from google.cloud import texttospeech as tts

class TextToSpeech:
    def __init__(
        self,
        gender='neutral',
        language='en-US',
        speed=1.0,
        effects=['handset-class-device']
    ):
        '''
        Create a TextToSpeech object. Parameters:
        -  gender:  'neutral' (default), 'female' or 'male'
        -  language code: 'en-US' (default), 'en-GB', 'cmn-CN', 'cmn-TW', ...
        -  speaking_rate: [0.25, 4.0] Default is 1.0
        -  effects_profile_id: list of effects
           (See https://cloud.google.com/text-to-speech/docs/audio-profiles)
        '''
        if gender == 'male':
            sex = tts.enums.SsmlVoiceGender.MALE
        elif gender == 'female':
            sex = tts.enums.SsmlVoiceGender.FEMALE
        else:
            sex == 'neutral'

        self.__client = tts.TextToSpeechClient()

        self.__voice  = tts.types.VoiceSelectionParams(
            language_code=language,
            ssml_gender=sex
        )
        self.__audio_config = tts.types.AudioConfig(
            audio_encoding=tts.enums.AudioEncoding.MP3,
            speaking_rate=speed,
            effects_profile_id=effects
        )

    def getMP3(self, text, output):
        if len(text) > 5000:
            print('** Speaker note is truncated > 5000 characters')

        input = tts.types.SynthesisInput(text=text)
        response = self.__client.synthesize_speech(
            input,
            self.__voice,
            self.__audio_config
        )
        with open(output, 'wb') as out:
            out.write(response.audio_content)
            print(f'Create "{output}"')

from pptx import Presentation
from lxml import etree
from pptx.util import Cm
import os

class Powerpoint:
    def __init__(self, pptx, tts=None):
        self.__fname = pptx
        self.__name  = os.path.splitext(pptx)[0]
        self.__pp    = Presentation(pptx)
        self.__tts   = tts

    def getSpkNotes(self):
        '''
        Retrieve a list of speaker notes
        '''
        sns = []
        for s in self.__pp.slides:
            if s.has_notes_slide:
                sns.append(s.notes_slide.notes_text_frame.text)
            else:
                sns.append('')
        else:
            return sns;

    def getSpkNotesGen(self):
        '''
        Generator to retrieve speaker notes
        '''
        for idx, s in enumerate(self.__pp.slides):
            if s.has_notes_slide:
                yield s.notes_slide.notes_text_frame.text
            else:
                yield None

    def VoiceAnnotatePP(self, out_pptx):
        '''
        Convert speaker notes to mp3 and insert into slides
        '''
        if not self.__tts:
            return

        for idx, s in enumerate(self.__pp.slides):
            print(f'Processing slide {idx+1} ... ', end='')
            if s.has_notes_slide:
                text=s.notes_slide.notes_text_frame.text
                outf = f'{self.__name}-{str(idx+1):>03}.mp3'
                self.__tts.getMP3(text, outf)
                movie = s.shapes.add_movie(
                    outf,
                    Cm(31.85), Cm(15.95), Cm(1.5), Cm(1.5),
                    poster_frame_image=None, mime_type='audio/mpeg3'
                )
                tree = movie._element.getparent().getparent().getnext().getnext().getnext()
                if tree is not None:
                    timing = [el for el in tree.iterdescendants() if etree.QName(el).localname == 'cond'][0]
                    timing.set('delay', '0')
                else:
                    tree = movie._element.getparent().getparent().getnext().getnext()
                    if tree is not None:
                        timing = [el for el in tree.iterdescendants() if etree.QName(el).localname == 'cond'][0]
                        timing.set('delay', '0')
            else:
                print('')

        self.__pp.save(out_pptx)

import argparse

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Text to Speech')
    parser.add_argument('source_pptx', type=str, help='Input PowerPoint file')
    parser.add_argument('output_pptx', type=str, help='Output PowerPoint file')
    parser.add_argument('--gender', nargs='?', default='female', help='Gender (default: female)')
    parser.add_argument('--lang',   nargs='?', default='en-GB',  help='Language (default: en-GB)')
    parser.add_argument('--speed',  nargs='?', default=1.0,      help='Speed (default: 1.0)', type=float)

    args = parser.parse_args()
    print(args.gender, args.lang, args.speed)

    t2s = TextToSpeech(gender=args.gender, language=args.lang, speed=args.speed)
    pp = Powerpoint(args.source_pptx, t2s)
    pp.VoiceAnnotatePP(args.output_pptx)
    print('Done')
